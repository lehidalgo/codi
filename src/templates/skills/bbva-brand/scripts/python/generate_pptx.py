#!/usr/bin/env python3
"""
generate_pptx.py — BBVA-branded PPTX generator (python-pptx, FALLBACK runtime).

Usage:
    python3 generate_pptx.py --content content.json --output output.pptx

Install dependency: pip install python-pptx
content.json schema: {title, subtitle?, author?, sections: [{number, label, heading, body, items?, callout?}]}
"""

import sys
import json
import argparse
from pathlib import Path
from datetime import datetime

sys.path.insert(0, str(Path(__file__).parent))
import brand_tokens as bt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(float(bt.LAYOUT["slide_width_in"]))
SLIDE_H = Inches(float(bt.LAYOUT["slide_height_in"]))
MARGIN = Inches(float(bt.LAYOUT["content_margin_in"]))
BAR_W = Inches(float(bt.LAYOUT["accent_bar_width_in"]))


def _rgb(key: str) -> RGBColor:
    r, g, b = bt.rgb(key)
    return RGBColor(r, g, b)


def _set_bg(slide, key: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(key)


def _add_rect(slide, x, y, w, h, color_key: str) -> None:
    shape = slide.shapes.add_shape(
        1,
        x,
        y,
        w,
        h,  # MSO_SHAPE_TYPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_key)
    shape.line.fill.background()


def _add_text(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    size: int,
    color_key: str,
    font_key: str = "pptx_body",
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
) -> None:
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = _rgb(color_key)
    run.font.name = bt.FONTS.get(font_key, "Arial")
    run.font.bold = bold
    run.font.italic = italic


def build_title_slide(prs: Presentation, content: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_bg(slide, "background_dark")
    _add_rect(slide, 0, 0, BAR_W, SLIDE_H, "primary")

    # BBVA wordmark
    _add_text(
        slide,
        "BBVA",
        MARGIN,
        Inches(0.3),
        Inches(2),
        Inches(0.5),
        22,
        "primary",
        "pptx_body",
        bold=True,
    )

    # Title
    _add_text(
        slide,
        content["title"],
        MARGIN,
        Inches(2.5),
        Inches(9),
        Inches(1.8),
        44,
        "white",
        "pptx_headlines",
        bold=True,
    )

    # Subtitle
    if content.get("subtitle"):
        _add_text(
            slide,
            content["subtitle"],
            MARGIN,
            Inches(4.5),
            Inches(9),
            Inches(0.8),
            20,
            "secondary",
            "pptx_body",
            italic=True,
        )

    # Footer
    author = content.get("author", "")
    year = str(datetime.now().year)
    footer = f"{author} · {year}" if author else year
    _add_text(
        slide,
        footer,
        MARGIN,
        Inches(6.8),
        Inches(12),
        Inches(0.4),
        11,
        "text_light",
        "pptx_body",
    )


def build_divider_slide(prs: Presentation, section: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, "primary")

    # Large section number (decorative)
    _add_text(
        slide,
        section["number"],
        Inches(0.5),
        Inches(1.5),
        Inches(4),
        Inches(2),
        96,
        "accent",
        "pptx_headlines",
        bold=True,
    )

    # Section heading
    _add_text(
        slide,
        section["heading"],
        Inches(0.5),
        Inches(3.5),
        Inches(10),
        Inches(1.5),
        32,
        "white",
        "pptx_headlines",
        bold=True,
    )

    # Section label
    _add_text(
        slide,
        section["label"].upper(),
        Inches(0.5),
        Inches(5.2),
        Inches(10),
        Inches(0.6),
        13,
        "secondary",
        "pptx_body",
    )


def build_content_slide(prs: Presentation, section: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, "background")
    _add_rect(slide, 0, 0, BAR_W, SLIDE_H, "primary")

    # Label
    _add_text(
        slide,
        section["label"].upper(),
        MARGIN,
        Inches(0.25),
        Inches(10),
        Inches(0.3),
        10,
        "primary",
        "pptx_body",
    )

    # Heading
    _add_text(
        slide,
        section["heading"],
        MARGIN,
        Inches(0.65),
        Inches(12),
        Inches(1),
        28,
        "primary_dark",
        "pptx_headlines",
        bold=True,
    )

    # Body
    _add_text(
        slide,
        section["body"],
        MARGIN,
        Inches(1.95),
        Inches(7.5),
        Inches(2.5),
        14,
        "text_primary",
        "pptx_body",
    )

    # Bullet items
    items = section.get("items", [])
    if items:
        txBox = slide.shapes.add_textbox(
            Inches(8.2), Inches(1.95), Inches(4.8), Inches(4.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = _rgb("text_primary")
            p.font.name = bt.FONTS.get("pptx_body", "Arial")

    # Callout
    if section.get("callout"):
        _add_rect(slide, MARGIN, Inches(4.7), Inches(7.5), Inches(1.5), "border")
        _add_text(
            slide,
            section["callout"],
            MARGIN + Inches(0.15),
            Inches(4.9),
            Inches(7.1),
            Inches(1.1),
            13,
            "primary_dark",
            "pptx_body",
            italic=True,
        )


def build_closing_slide(prs: Presentation, content: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, "background_dark")
    _add_rect(slide, 0, 0, BAR_W, SLIDE_H, "accent")

    _add_text(
        slide,
        "Gracias",
        MARGIN,
        Inches(2.5),
        Inches(8),
        Inches(1.5),
        52,
        "white",
        "pptx_headlines",
        bold=True,
    )
    _add_text(
        slide,
        content["title"],
        MARGIN,
        Inches(4.2),
        Inches(10),
        Inches(0.6),
        16,
        "secondary",
        "pptx_body",
    )
    _add_text(
        slide,
        "BBVA",
        MARGIN,
        Inches(6.5),
        Inches(3),
        Inches(0.5),
        18,
        "primary",
        "pptx_body",
        bold=True,
    )


def generate_pptx(content: dict, output_path: str) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs, content)
    for section in content.get("sections", []):
        build_divider_slide(prs, section)
        build_content_slide(prs, section)
    build_closing_slide(prs, content)

    prs.save(output_path)
    n = len(content.get("sections", []))
    print(f"PPTX written: {output_path} ({n} sections)")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="BBVA PPTX generator (python-pptx fallback)"
    )
    parser.add_argument("--content", required=True, help="Path to content.json")
    parser.add_argument("--output", default="output.pptx", help="Output .pptx file")
    args = parser.parse_args()

    with open(args.content, "r", encoding="utf-8") as f:
        content = json.load(f)

    generate_pptx(content, args.output)


if __name__ == "__main__":
    main()

"""
generate_pptx.py — Generate an RL3-branded PowerPoint presentation from structured content.

CLI: python generate_pptx.py --content content.json --output output.pptx
"""

import sys
import os
import argparse
import json
import zipfile
import io
import re
import tempfile

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import brand_tokens as bt

# ── Brand constants ────────────────────────────────────────────────

BG_COLOR = RGBColor(0x0A, 0x0A, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC8, 0xB8, 0x8A)
GRAY = RGBColor(0x7A, 0x7A, 0x7A)

FONT_HEADLINES = bt.FONTS["headlines"]   # Space Grotesk
FONT_MONO = bt.FONTS["mono"]            # Space Mono
FONT_BODY = bt.FONTS["body"]            # Instrument Sans

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FIXED_TIMESTAMP = "2026-01-01T00:00:00Z"


# ── Helpers ────────────────────────────────────────────────────────

def _set_slide_bg(slide):
    """Apply the standard RL3 dark background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def _add_textbox(slide, left, top, width, height):
    """Add a textbox and return it."""
    return slide.shapes.add_textbox(left, top, width, height)


def _add_run(paragraph, text, font_name, font_size, color, bold=False, italic=False):
    """Add a formatted run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return run


# ── Validation ─────────────────────────────────────────────────────

def _validate_content(content):
    """Validate required fields, raise ValueError if missing."""
    if "title" not in content:
        raise ValueError("Missing required field: 'title'")
    if "sections" not in content:
        raise ValueError("Missing required field: 'sections'")
    if not isinstance(content["sections"], list) or len(content["sections"]) < 1:
        raise ValueError("'sections' must be a non-empty list")
    for i, section in enumerate(content["sections"]):
        for field in ("number", "label", "heading", "body"):
            if field not in section:
                raise ValueError(f"Section {i}: missing required field '{field}'")


# ── Slide builders ─────────────────────────────────────────────────

def _build_title_slide(prs, content):
    """Build the title slide with RL3 logo and subtitle."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    # Logo: "RL" white + "3" gold
    logo_box = _add_textbox(
        slide,
        Inches(4.5), Inches(2.0),
        Inches(4.333), Inches(1.5),
    )
    tf = logo_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, "RL", FONT_HEADLINES, Pt(72), WHITE, bold=True)
    _add_run(p, "3", FONT_HEADLINES, Pt(72), GOLD, bold=True)

    # "AI AGENCY" subtitle
    agency_box = _add_textbox(
        slide,
        Inches(4.5), Inches(3.5),
        Inches(4.333), Inches(0.6),
    )
    tf2 = agency_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    _add_run(p2, "AI AGENCY", FONT_MONO, Pt(14), GRAY)

    # Content subtitle if present
    if content.get("subtitle"):
        sub_box = _add_textbox(
            slide,
            Inches(3.5), Inches(4.3),
            Inches(6.333), Inches(0.6),
        )
        tf3 = sub_box.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        _add_run(p3, content["subtitle"], FONT_BODY, Pt(18), WHITE)

    return slide


def _build_divider_slide(prs, section):
    """Build a section divider slide with large number and heading."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    # Large section number
    num_box = _add_textbox(
        slide,
        Inches(1.5), Inches(2.0),
        Inches(3.0), Inches(2.0),
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    _add_run(p, section["number"], FONT_HEADLINES, Pt(96), GOLD, bold=True)

    # Section label
    label_box = _add_textbox(
        slide,
        Inches(1.5), Inches(4.0),
        Inches(6.0), Inches(0.5),
    )
    tf2 = label_box.text_frame
    p2 = tf2.paragraphs[0]
    _add_run(p2, section["label"].upper(), FONT_MONO, Pt(12), GRAY)

    # Section heading
    heading_box = _add_textbox(
        slide,
        Inches(1.5), Inches(4.6),
        Inches(10.0), Inches(1.2),
    )
    tf3 = heading_box.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    _add_run(p3, section["heading"], FONT_HEADLINES, Pt(36), WHITE, bold=True)

    return slide


def _build_content_slide(prs, section):
    """Build a section content slide with body, items, and optional callout."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    # Body text
    body_box = _add_textbox(
        slide,
        Inches(1.5), Inches(1.5),
        Inches(10.0), Inches(1.5),
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_run(p, section["body"], FONT_BODY, Pt(18), WHITE)

    # Items as bullet list with arrow prefix
    items = section.get("items", [])
    if items:
        items_box = _add_textbox(
            slide,
            Inches(1.5), Inches(3.2),
            Inches(10.0), Inches(2.5),
        )
        tf2 = items_box.text_frame
        tf2.word_wrap = True
        for idx, item in enumerate(items):
            if idx == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.space_after = Pt(6)
            _add_run(p2, f"\u2192 {item}", FONT_BODY, Pt(16), WHITE)

    # Callout if present
    if section.get("callout"):
        callout_box = _add_textbox(
            slide,
            Inches(1.5), Inches(5.5),
            Inches(10.0), Inches(0.8),
        )
        tf3 = callout_box.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        _add_run(p3, section["callout"], FONT_BODY, Pt(16), GOLD, italic=True)

    return slide


def _build_final_slide(prs):
    """Build the closing slide with the RL3 cycle and tagline."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    # Cycle label
    cycle_box = _add_textbox(
        slide,
        Inches(3.0), Inches(2.8),
        Inches(7.333), Inches(1.0),
    )
    tf = cycle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, bt.CYCLE_LABEL, FONT_HEADLINES, Pt(36), WHITE, bold=True)

    # Tagline
    tag_box = _add_textbox(
        slide,
        Inches(3.0), Inches(4.0),
        Inches(7.333), Inches(0.8),
    )
    tf2 = tag_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    _add_run(p2, bt.TAGLINE, FONT_BODY, Pt(16), GRAY, italic=True)

    return slide


# ── Determinism: strip timestamps from pptx ZIP ───────────────────

def _strip_timestamps(pptx_path):
    """
    Re-write the pptx ZIP so that docProps/core.xml has fixed timestamps
    and all ZIP entries have a consistent date, ensuring SHA-256 determinism.
    """
    buf = io.BytesIO()
    with open(pptx_path, "rb") as f:
        buf.write(f.read())
    buf.seek(0)

    out_buf = io.BytesIO()
    fixed_date = (2026, 1, 1, 0, 0, 0)

    with zipfile.ZipFile(buf, "r") as zin:
        with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)

                # Fix timestamps in core.xml
                if item.filename == "docProps/core.xml":
                    text = data.decode("utf-8")
                    # Replace dcterms:created and dcterms:modified values
                    text = re.sub(
                        r"(<dcterms:created[^>]*>)[^<]*(</dcterms:created>)",
                        rf"\g<1>{FIXED_TIMESTAMP}\g<2>",
                        text,
                    )
                    text = re.sub(
                        r"(<dcterms:modified[^>]*>)[^<]*(</dcterms:modified>)",
                        rf"\g<1>{FIXED_TIMESTAMP}\g<2>",
                        text,
                    )
                    data = text.encode("utf-8")

                # Set consistent date_time on every entry
                info = zipfile.ZipInfo(item.filename, date_time=fixed_date)
                info.compress_type = zipfile.ZIP_DEFLATED
                # Preserve external attributes (directory flag etc.)
                info.external_attr = item.external_attr
                zout.writestr(info, data)

    with open(pptx_path, "wb") as f:
        f.write(out_buf.getvalue())


# ── Public API ─────────────────────────────────────────────────────

def generate_pptx(content: dict, output_path: str):
    """
    Generate an RL3-branded PPTX presentation from structured content.

    Args:
        content: dict with keys title, sections (list), optional subtitle/date/footer_text
        output_path: filesystem path for the resulting .pptx file
    """
    _validate_content(content)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. Title slide
    _build_title_slide(prs, content)

    # 2. Per section: divider + content slides
    for section in content["sections"]:
        _build_divider_slide(prs, section)
        _build_content_slide(prs, section)

    # 3. Final slide
    _build_final_slide(prs)

    # Save
    prs.save(output_path)

    # Strip timestamps for determinism
    _strip_timestamps(output_path)


# ── CLI ────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Generate an RL3-branded PPTX from a JSON content file."
    )
    parser.add_argument("--content", required=True, help="Path to JSON content file")
    parser.add_argument("--output", required=True, help="Output .pptx file path")
    args = parser.parse_args()

    if not os.path.isfile(args.content):
        print(f"Error: content file not found: {args.content}", file=sys.stderr)
        sys.exit(1)

    with open(args.content, "r", encoding="utf-8") as f:
        content = json.load(f)

    generate_pptx(content, args.output)
    print(f"Generated: {args.output}")


if __name__ == "__main__":
    main()

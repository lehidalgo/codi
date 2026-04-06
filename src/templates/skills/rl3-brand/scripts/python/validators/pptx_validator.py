"""
pptx_validator.py — Validate a PowerPoint file against RL3 brand guidelines.

Rules:
  1. slide_background   — All slides have solid fill #0a0a0b
  2. logo_3_color       — Where "3" text appears, its run color is #c8b88a (gold)
  3. title_subtitle     — First slide contains "AI AGENCY" text somewhere
  4. section_divider    — At least one slide with a text run >= Pt(48) in gold
  5. forbidden_phrases  — No phrase from brand_tokens.PHRASES_AVOID in any text frame

CLI: python pptx_validator.py --input file.pptx
"""

import sys
import os
import argparse
import json

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

import brand_tokens as bt

# ── Expected constants ──────────────────────────────────────────────

BG_COLOR = RGBColor(0x0A, 0x0A, 0x0B)
GOLD_COLOR = RGBColor(0xC8, 0xB8, 0x8A)
MIN_DIVIDER_SIZE = Pt(48)


# ── Rule checkers ───────────────────────────────────────────────────


def _check_slide_background(prs):
    """Rule 1: every slide must have a solid fill background of #0a0a0b."""
    errors = []
    for idx, slide in enumerate(prs.slides, start=1):
        bg = slide.background
        fill = bg.fill
        # Check if fill is solid
        if fill.type is None:
            errors.append(
                {
                    "rule": "slide_background",
                    "message": f"Slide {idx}: no solid background fill set.",
                }
            )
            continue

        # fill.type == MSO_FILL.SOLID when solid() was called
        try:
            color = fill.fore_color.rgb
        except Exception:
            errors.append(
                {
                    "rule": "slide_background",
                    "message": f"Slide {idx}: background fill is not a solid RGB colour.",
                }
            )
            continue
        if color != BG_COLOR:
            errors.append(
                {
                    "rule": "slide_background",
                    "message": f"Slide {idx}: background is #{color}, expected #{BG_COLOR}.",
                }
            )
    return errors


def _check_logo_3_color(prs):
    """Rule 2: in the RL3 logo, the '3' run must be gold.

    The logo is identified by a paragraph whose combined text contains 'RL3'
    (or starts with 'RL' followed by a run whose text is '3'). Only those
    runs are checked — content text like 'Fase 3' is ignored.
    """
    errors = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in para.runs)
                if "RL" not in full_text:
                    continue
                # This paragraph contains an RL3 logo candidate
                for run in para.runs:
                    if run.text.strip() == "3":
                        try:
                            color = run.font.color.rgb
                        except Exception:
                            color = None
                        if color != GOLD_COLOR:
                            errors.append(
                                {
                                    "rule": "logo_3_color",
                                    "message": (
                                        f"Slide {idx}: logo '3' run "
                                        f"has colour #{color}, expected #{GOLD_COLOR}."
                                    ),
                                }
                            )
    return errors


def _check_title_subtitle(prs):
    """Rule 3: first slide must contain 'AI AGENCY' somewhere in its text."""
    errors = []
    if len(prs.slides) == 0:
        errors.append(
            {
                "rule": "title_subtitle",
                "message": "Presentation has no slides.",
            }
        )
        return errors

    first_slide = prs.slides[0]
    found = False
    for shape in first_slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in para.runs)
            if "AI AGENCY" in full_text:
                found = True
                break
        if found:
            break
    if not found:
        # Also check shape.text as fallback
        for shape in first_slide.shapes:
            if not shape.has_text_frame:
                continue
            if "AI AGENCY" in shape.text_frame.text:
                found = True
                break
    if not found:
        errors.append(
            {
                "rule": "title_subtitle",
                "message": "First slide does not contain 'AI AGENCY' text.",
            }
        )
    return errors


def _check_section_divider(prs):
    """Rule 4: at least one slide must have a text run >= Pt(48) in gold colour."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    size = run.font.size
                    if size is not None and size >= MIN_DIVIDER_SIZE:
                        try:
                            color = run.font.color.rgb
                        except Exception:
                            continue
                        if color == GOLD_COLOR:
                            return []  # rule passes
    return [
        {
            "rule": "section_divider",
            "message": "No slide has a text run >= 48pt in gold (#c8b88a).",
        }
    ]


def _check_forbidden_phrases(prs):
    """Rule 5: no text frame may contain any phrase from PHRASES_AVOID."""
    errors = []
    forbidden = [p.lower() for p in bt.PHRASES_AVOID]
    forbidden_orig = bt.PHRASES_AVOID

    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            frame_text = shape.text_frame.text.lower()
            for i, phrase in enumerate(forbidden):
                if phrase in frame_text:
                    errors.append(
                        {
                            "rule": "forbidden_phrases",
                            "message": (
                                f"Slide {idx}: forbidden phrase "
                                f"'{forbidden_orig[i]}' found."
                            ),
                        }
                    )
    return errors


# ── Public API ──────────────────────────────────────────────────────


def validate_pptx(filepath: str) -> dict:
    """
    Validate a .pptx file against RL3 brand rules.

    Returns:
        {"passed": bool, "errors": [...], "warnings": []}
    """
    prs = Presentation(filepath)

    errors = []
    errors.extend(_check_slide_background(prs))
    errors.extend(_check_logo_3_color(prs))
    errors.extend(_check_title_subtitle(prs))
    errors.extend(_check_section_divider(prs))
    errors.extend(_check_forbidden_phrases(prs))

    return {
        "passed": len(errors) == 0,
        "errors": errors,
        "warnings": [],
    }


# ── CLI ─────────────────────────────────────────────────────────────


def main():
    parser = argparse.ArgumentParser(
        description="Validate a PPTX file against RL3 brand guidelines."
    )
    parser.add_argument("--input", required=True, help="Path to .pptx file")
    args = parser.parse_args()

    if not os.path.isfile(args.input):
        print(f"Error: file not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    result = validate_pptx(args.input)
    print(json.dumps(result, indent=2))
    sys.exit(0 if result["passed"] else 1)


if __name__ == "__main__":
    main()

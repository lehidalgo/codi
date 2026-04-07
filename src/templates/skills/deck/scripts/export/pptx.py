#!/usr/bin/env python3
"""
Codi Deck — PPTX Export
Usage: python3 pptx.py --input deck.html --output deck.pptx

Parses slide HTML (title, headings, lists, paragraphs) and builds a .pptx file
using python-pptx. Requires: pip install python-pptx beautifulsoup4 lxml

Each .slide element becomes one PPTX slide (16:9 widescreen, 960x540 px).
"""

import argparse
import os
import sys

try:
    from bs4 import BeautifulSoup
except ImportError:
    sys.exit("beautifulsoup4 not installed. Run: pip install beautifulsoup4 lxml")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx")

# ========== Constants ==========

SLIDE_W_IN = 10.0  # inches  (960 / 96 dpi)
SLIDE_H_IN = 5.625  # inches  (540 / 96 dpi)
MARGIN_IN = 0.6
CONTENT_W = SLIDE_W_IN - MARGIN_IN * 2
CONTENT_H = SLIDE_H_IN - MARGIN_IN * 2

COLOR_PRIMARY = RGBColor(0x00, 0x13, 0x91)
COLOR_TEXT = RGBColor(0x1A, 0x1A, 0x2A)
COLOR_SECONDARY = RGBColor(0x4A, 0x4A, 0x68)
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ACCENT = RGBColor(0x85, 0xC8, 0xFF)

FONT_HEADLINE = "Source Serif 4"
FONT_BODY = "Lato"

# ========== Helpers ==========


def parse_color(hex_str: str) -> RGBColor | None:
    """Parse a CSS hex color like #001391 or 001391."""
    s = hex_str.strip().lstrip("#")
    if len(s) == 6:
        try:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
        except ValueError:
            pass
    return None


def is_accent_slide(el) -> bool:
    """True if the slide element has the --accent modifier class."""
    return "slide--accent" in el.get("class", [])


def add_text_box(
    slide,
    text: str,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    font_name: str,
    font_size_pt: float,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_slide_background(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ========== Slide Builders ==========


def build_slide(prs: Presentation, el) -> None:
    """Convert one <section class="slide"> element to a PPTX slide."""
    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)

    accent = is_accent_slide(el)
    bg_color = COLOR_PRIMARY if accent else COLOR_BG
    text_color = RGBColor(0xFF, 0xFF, 0xFF) if accent else COLOR_TEXT
    subtext_color = RGBColor(0xCC, 0xD6, 0xFF) if accent else COLOR_SECONDARY

    set_slide_background(slide, bg_color)

    y = MARGIN_IN

    # --- Title (.title--xl, .title--lg, h1) ---
    title_el = (
        el.select_one(".title--xl")
        or el.select_one(".title--lg")
        or el.select_one("h1")
    )
    if title_el:
        title_text = title_el.get_text(strip=True)
        font_size = 40.0 if ".title--xl" in el.decode() else 32.0
        add_text_box(
            slide,
            title_text,
            MARGIN_IN,
            y,
            CONTENT_W,
            1.4,
            FONT_HEADLINE,
            font_size,
            bold=True,
            color=text_color,
            align=PP_ALIGN.CENTER
            if "centered" in str(el.get("class", []))
            else PP_ALIGN.LEFT,
        )
        y += 1.5

    # --- Section number (.section-number) ---
    num_el = el.select_one(".section-number")
    if num_el and not title_el:
        add_text_box(
            slide,
            num_el.get_text(strip=True),
            MARGIN_IN,
            y,
            CONTENT_W,
            1.0,
            FONT_BODY,
            60.0,
            bold=True,
            color=COLOR_ACCENT,
            align=PP_ALIGN.LEFT,
        )
        y += 1.1

    # --- Subtitle (.subtitle) ---
    sub_el = el.select_one(".subtitle")
    if sub_el:
        add_text_box(
            slide,
            sub_el.get_text(strip=True),
            MARGIN_IN,
            y,
            CONTENT_W,
            0.5,
            FONT_BODY,
            18.0,
            color=subtext_color,
            align=PP_ALIGN.CENTER
            if "centered" in str(el.get("class", []))
            else PP_ALIGN.LEFT,
        )
        y += 0.6

    # --- Meta (.meta) ---
    meta_el = el.select_one(".meta")
    if meta_el:
        add_text_box(
            slide,
            meta_el.get_text(strip=True),
            MARGIN_IN,
            y,
            CONTENT_W,
            0.4,
            FONT_BODY,
            13.0,
            color=subtext_color,
            align=PP_ALIGN.CENTER
            if "centered" in str(el.get("class", []))
            else PP_ALIGN.LEFT,
        )
        y += 0.5

    # --- H2 heading ---
    h2_el = el.select_one("h2")
    if h2_el and not title_el:
        add_text_box(
            slide,
            h2_el.get_text(strip=True),
            MARGIN_IN,
            y,
            CONTENT_W,
            0.9,
            FONT_HEADLINE,
            28.0,
            bold=True,
            color=text_color,
        )
        y += 1.0

    # --- Bullet list (.bullet-list) ---
    list_el = el.select_one(".bullet-list")
    if list_el:
        items = [li.get_text(strip=True) for li in list_el.find_all("li")]
        for item in items:
            if y + 0.4 > SLIDE_H_IN - MARGIN_IN:
                break
            txBox = slide.shapes.add_textbox(
                Inches(MARGIN_IN + 0.2),
                Inches(y),
                Inches(CONTENT_W - 0.2),
                Inches(0.38),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "• " + item
            run.font.name = FONT_BODY
            run.font.size = Pt(16)
            run.font.color.rgb = text_color
            y += 0.4

    # --- Body paragraph (p) ---
    body_p = el.select_one(".slide__content > p")
    if body_p:
        add_text_box(
            slide,
            body_p.get_text(strip=True),
            MARGIN_IN,
            y,
            CONTENT_W,
            0.9,
            FONT_BODY,
            15.0,
            color=text_color,
        )


# ========== Main ==========


def main() -> None:
    parser = argparse.ArgumentParser(description="Export codi deck HTML to PPTX")
    parser.add_argument("--input", default="deck.html", help="Input HTML file")
    parser.add_argument("--output", default="deck.pptx", help="Output PPTX file")
    args = parser.parse_args()

    abs_input = os.path.abspath(args.input)
    abs_output = os.path.abspath(args.output)

    if not os.path.exists(abs_input):
        sys.exit(f"Input file not found: {abs_input}")

    with open(abs_input, encoding="utf-8") as fh:
        soup = BeautifulSoup(fh, "lxml")

    slide_els = soup.select(".slide")
    if not slide_els:
        sys.exit(f"No .slide elements found in {abs_input}")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    for el in slide_els:
        build_slide(prs, el)

    os.makedirs(os.path.dirname(abs_output) or ".", exist_ok=True)
    prs.save(abs_output)
    print(f"PPTX exported: {abs_output} ({len(slide_els)} slides)")


if __name__ == "__main__":
    main()
